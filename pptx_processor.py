from pptx import Presentation
from pptx.util import Pt, Emu
import os
import re
import traceback

class PPTXProcessor:
    def __init__(self, translator, translation_level="normal"):
        self.translator = translator
        self.korean_font = "Malgun Gothic"
        self.translation_level = translation_level
        self.slide_height = None  # Set during process_presentation

    def _should_skip_shape(self, shape):
        """Determine if a shape's text should be excluded from translation."""
        level = self.translation_level

        # 1. Title/Subtitle placeholders - skip for normal and minimal
        if level in ("normal", "minimal") and shape.is_placeholder:
            try:
                idx = shape.placeholder_format.idx
                # 0=title, 1=center title, 13=subtitle
                if idx in (0, 1, 13):
                    return True
            except Exception:
                pass
            try:
                from pptx.enum.shapes import PP_PLACEHOLDER
                ptype = shape.placeholder_format.type
                if ptype in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.SUBTITLE,
                             PP_PLACEHOLDER.CENTER_TITLE):
                    return True
            except Exception:
                pass

        # 2. Large font heuristic (diagram/chart display labels)
        #    GUARD: Only apply to shapes with few paragraphs (≤3).
        #    Content-rich text boxes (e.g., "Business Challenges" heading + 5 bullets)
        #    should NOT be skipped just because their heading has a large font.
        if shape.has_text_frame:
            non_empty_for_font = [p for p in shape.text_frame.paragraphs if p.text.strip()]
            if len(non_empty_for_font) <= 3:
                if level == "minimal":
                    font_threshold = Pt(20)
                elif level == "normal":
                    font_threshold = Pt(24)
                else:  # thorough
                    font_threshold = Pt(32)

                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size and run.font.size >= font_threshold:
                            return True

        # 3. Footer/watermark detection - shape near bottom of slide with short text
        if level in ("normal", "minimal") and shape.has_text_frame and self.slide_height:
            try:
                shape_top = shape.top
                if shape_top is not None and self.slide_height > 0:
                    # Shape is in the bottom 15% of the slide
                    if shape_top > self.slide_height * 0.85:
                        text = shape.text_frame.text.strip()
                        if len(text) <= 30:
                            return True
            except Exception:
                pass

        # 4. Label-only shape detection - short text, no sentence structure, few paragraphs
        #    Catches: diagram labels, process step names, small callout shapes
        if level in ("normal", "minimal") and shape.has_text_frame:
            tf = shape.text_frame
            non_empty_paras = [p for p in tf.paragraphs if p.text.strip()]
            total_text = tf.text.strip()

            if non_empty_paras and len(non_empty_paras) <= 3:
                words = total_text.split()
                has_sentence_punct = bool(re.search(r'[.!?。]', total_text))

                if level == "minimal":
                    max_chars, max_words = 80, 10
                else:  # normal
                    max_chars, max_words = 60, 8

                if (len(total_text) <= max_chars
                        and len(words) <= max_words
                        and not has_sentence_punct):
                    return True

        # 5. Short label detection (very short non-sentence text) - original rule
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            max_len = 8 if level == "minimal" else 5
            if 1 < len(text) <= max_len and ' ' not in text:
                return True

        return False

    def _should_skip_text(self, text):
        """Check if a specific text string should skip translation."""
        stripped = text.strip()
        if not stripped:
            return True

        # Already Korean
        if any('\uac00' <= c <= '\ud7a3' for c in stripped):
            return True

        # Pure numbers, dates, or version strings
        if re.match(r'^[\d\.\-/\s:,]+$', stripped):
            return True

        # All-caps short acronyms (2-6 chars, allowing /)
        if re.match(r'^[A-Z/]{2,6}$', stripped):
            return True

        # Fiori app IDs (e.g., "F1234", "F0842")
        if re.match(r'^F\d{4,5}$', stripped):
            return True

        # Check against do-not-translate glossary entries
        if hasattr(self.translator, '_do_not_translate'):
            if stripped.lower() in self.translator._do_not_translate:
                return True

        return False

    def _is_heading_paragraph(self, paragraph, text_frame):
        """Detect if a paragraph is a category heading at the top of a content text box.

        Pattern: first paragraph is short (e.g., "Business Challenges"),
        followed by multiple paragraphs of body/bullet content.
        """
        level = self.translation_level
        if level == "thorough":
            return False

        text = paragraph.text.strip()
        if not text:
            return False

        # Get all non-empty paragraphs in this text frame
        all_paras = [p for p in text_frame.paragraphs if p.text.strip()]

        # Must have at least 3 paragraphs total (heading + 2+ body items)
        if len(all_paras) < 3:
            return False

        # Only applies to the FIRST non-empty paragraph
        if all_paras[0] is not paragraph:
            return False

        # First paragraph must be short and not look like a sentence
        if len(text) > 40:
            return False
        if re.search(r'[.!?。;]$', text):
            return False

        return True

    def process_presentation(self, input_path, output_path, progress_callback=None):
        try:
            from concurrent.futures import ThreadPoolExecutor
            prs = Presentation(input_path)

            # Store slide height for footer/watermark detection
            self.slide_height = prs.slide_height

            # Step 1: Collect all text targets
            text_frames = []
            def collect_frames(shapes, is_master=False, nesting_depth=0):
                if not shapes: return
                for shape in shapes:
                    # For Masters/Layouts, only translate if it's NOT a placeholder
                    # and has actual content (placeholders often overlap or cause shifts)
                    if is_master and (shape.is_placeholder or not shape.has_text_frame):
                        continue

                    # Smart shape-level filtering (not applied to master slides)
                    if not is_master and self._should_skip_shape(shape):
                        continue

                    if shape.has_text_frame:
                        # For deeply nested group shapes (depth >= 2),
                        # skip short text as it's likely a diagram node label
                        if nesting_depth >= 2:
                            text = shape.text_frame.text.strip()
                            if len(text) < 30:
                                continue
                        text_frames.append(shape.text_frame)
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                text_frames.append(cell.text_frame)
                    if shape.shape_type == 6:  # Group shape
                        collect_frames(shape.shapes, is_master, nesting_depth + 1)

            # 1.1: Collect from slides (Higher priority)
            for slide in prs.slides:
                collect_frames(slide.shapes, is_master=False)
                if slide.has_notes_slide:
                    collect_frames(slide.notes_slide.shapes, is_master=False)

            # 1.2: Collect from layouts/masters (Only if user wanted 'thorough' but stay safe)
            for master in prs.slide_masters:
                collect_frames(master.shapes, is_master=True)
                for layout in master.slide_layouts:
                    collect_frames(layout.shapes, is_master=True)

            # Step 2: Extract all unique paragraphs (with text-level and heading filtering)
            unique_texts = set()
            paragraphs_to_translate = []
            for tf in text_frames:
                for p in tf.paragraphs:
                    if p.text.strip() and len(p.text.strip()) > 1: # Skip single chars/bullets
                        # Skip category headings (e.g., "Business Challenges" at top of text box)
                        if self._is_heading_paragraph(p, tf):
                            continue
                        if not self._should_skip_text(p.text):
                            unique_texts.add(p.text)
                            paragraphs_to_translate.append(p)

            # Step 3: Translate unique texts in parallel
            total_unique = len(unique_texts)
            translation_map = {}
            translation_errors = []
            processed_count = 0

            # Only translate if there's actually something to translate
            if unique_texts:
                max_workers = 15
                with ThreadPoolExecutor(max_workers=max_workers) as executor:
                    future_to_text = {executor.submit(self.translator.translate, text): text for text in unique_texts}

                    for i, future in enumerate(future_to_text):
                        text = future_to_text[future]
                        try:
                            result = future.result()
                            if result:
                                translation_map[text] = result
                            else:
                                translation_map[text] = text
                        except Exception as e:
                            err_msg = f"'{text[:30]}...' 번역 실패 ({type(e).__name__}): {str(e)}"
                            translation_errors.append(err_msg)
                            translation_map[text] = text

                        processed_count += 1
                        if progress_callback:
                            progress_callback((processed_count / total_unique) * 0.8)

            # Step 4: Apply translations
            for i, p in enumerate(paragraphs_to_translate):
                if p.text in translation_map:
                    translated_text = translation_map[p.text]
                    if translated_text and translated_text != p.text:
                        self._update_paragraph_text(p, translated_text)

                if progress_callback:
                    progress_callback(0.8 + (i / len(paragraphs_to_translate)) * 0.2)

            prs.save(output_path)
            return output_path, translation_errors
        except Exception as e:
            print(f"Critical error in process_presentation: {e}")
            traceback.print_exc()
            raise e

    def _update_paragraph_text(self, paragraph, translated_text):
        """Update paragraph text while preserving ALL original formatting.

        Key insight: Instead of paragraph.text = ... (which destroys all runs/XML),
        we modify individual run.text in-place so all XML attributes are preserved.
        """
        if translated_text is None:
            return

        try:
            runs = paragraph.runs

            if not runs:
                # No runs exist — fallback to direct assignment (rare edge case)
                paragraph.text = translated_text
                return

            if len(runs) == 1:
                # Single run: simply replace its text (keeps ALL formatting)
                runs[0].text = translated_text
            else:
                # Multiple runs: put all translated text in the first run,
                # clear the remaining runs (preserves first run's formatting)
                runs[0].text = translated_text
                for run in runs[1:]:
                    run.text = ""

            # Only set Korean font name — do NOT touch size, color, bold, etc.
            # They are already preserved from the original run XML
            for run in runs:
                if run.text:  # Only set font on runs that have text
                    run.font.name = self.korean_font

        except Exception as e:
            print(f"Error updating paragraph text: {e}")
            traceback.print_exc()
