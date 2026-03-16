import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
import io
import os
import re
import traceback


class PDFProcessor:
    """Convert PDF to PPTX with page images as backgrounds and text boxes overlaid.

    The resulting PPTX can then be translated by the existing PPTXProcessor.
    Uses OCR (Tesseract) as fallback for pages with image-embedded text.
    """

    PT_TO_EMU = 12700  # 1 PDF pt = 12700 EMU

    # Minimum non-footer text blocks to consider a page "has text"
    OCR_THRESHOLD = 3

    def __init__(self):
        # Ensure Tesseract is in PATH
        tesseract_path = r"C:\Program Files\Tesseract-OCR"
        if os.path.exists(tesseract_path) and tesseract_path not in os.environ.get("PATH", ""):
            os.environ["PATH"] = tesseract_path + ";" + os.environ["PATH"]

        self._ocr_available = self._check_ocr()

    def _check_ocr(self):
        """Check if Tesseract OCR is available."""
        try:
            import subprocess
            result = subprocess.run(
                ["tesseract", "--version"],
                capture_output=True, text=True, timeout=5
            )
            return result.returncode == 0
        except Exception:
            return False

    def _extract_line_text_with_spacing(self, line):
        """Extract text from a line's spans with proper spacing."""
        spans = line.get("spans", [])
        if not spans:
            return "", 0, False

        parts = []
        prev_end_x = None
        max_size = 0
        is_bold = False

        for span in spans:
            text = span["text"]
            if not text:
                continue

            if prev_end_x is not None:
                gap = span["origin"][0] - prev_end_x
                char_width = span["size"] * 0.5
                if gap > char_width * 0.3:
                    parts.append(" ")

            parts.append(text)
            prev_end_x = span.get("bbox", [0, 0, 0, 0])[2]

            if span["size"] > max_size:
                max_size = span["size"]
            if "bold" in span.get("font", "").lower():
                is_bold = True

        return "".join(parts), max_size, is_bold

    def _is_footer_or_header(self, text, y_top, y_bottom, page_height):
        """Detect footer/header/copyright text."""
        stripped = text.strip()

        if re.match(r'^\d+\s*INTERNAL', stripped, re.IGNORECASE):
            return True
        if re.search(r'INTERNAL\s*[–\-]\s*SAP', stripped, re.IGNORECASE):
            return True
        if y_bottom > page_height * 0.90 and re.search(r'INTERNAL', stripped, re.IGNORECASE):
            return True
        if '© SAP' in stripped or 'All rights reserved' in stripped:
            return True
        if y_bottom > page_height * 0.90 and len(stripped) <= 30:
            return True

        return False

    def _sample_bg_color(self, pixmap, bbox, zoom):
        """Sample the dominant background color around a text area."""
        x0 = max(0, min(int(bbox[0] * zoom), pixmap.width - 1))
        y0 = max(0, min(int(bbox[1] * zoom), pixmap.height - 1))
        x1 = max(0, min(int(bbox[2] * zoom), pixmap.width - 1))
        y1 = max(0, min(int(bbox[3] * zoom), pixmap.height - 1))

        sample_points = []
        mid_y = (y0 + y1) // 2
        mid_x = (x0 + x1) // 2
        if x0 > 2:
            sample_points.append((x0 - 2, mid_y))
        if y0 > 2:
            sample_points.append((mid_x, y0 - 2))
        if x1 < pixmap.width - 2:
            sample_points.append((x1 + 1, mid_y))

        if not sample_points:
            return RGBColor(255, 255, 255)

        r_total, g_total, b_total = 0, 0, 0
        valid = 0
        for px, py in sample_points:
            try:
                pixel = pixmap.pixel(px, py)
                r_total += pixel[0]
                g_total += pixel[1]
                b_total += pixel[2]
                valid += 1
            except Exception:
                continue

        if valid == 0:
            return RGBColor(255, 255, 255)

        return RGBColor(r_total // valid, g_total // valid, b_total // valid)

    def _get_font_color(self, line):
        """Extract the dominant font color from a line's spans."""
        for span in line.get("spans", []):
            color = span.get("color", 0)
            if isinstance(color, int):
                r = (color >> 16) & 0xFF
                g = (color >> 8) & 0xFF
                b = color & 0xFF
                return RGBColor(r, g, b)
        return None

    def _count_real_text_blocks(self, text_dict, page_height):
        """Count text blocks excluding footers/headers."""
        count = 0
        for block in text_dict.get("blocks", []):
            if block["type"] != 0:
                continue
            for line in block.get("lines", []):
                line_text, _, _ = self._extract_line_text_with_spacing(line)
                line_text = line_text.strip()
                if not line_text:
                    continue
                bbox = line.get("bbox", block["bbox"])
                if not self._is_footer_or_header(line_text, bbox[1], bbox[3], page_height):
                    count += 1
        return count

    def _add_text_boxes_from_dict(self, slide, text_dict, page_height, pix, zoom,
                                   slide_width_emu, slide_height_emu):
        """Create text boxes from a text dict (normal extraction or OCR)."""
        count = 0
        for block in text_dict.get("blocks", []):
            if block["type"] != 0:
                continue

            for line in block.get("lines", []):
                line_text, font_size, is_bold = self._extract_line_text_with_spacing(line)
                line_text = line_text.strip()
                if not line_text:
                    continue

                line_bbox = line.get("bbox", block["bbox"])
                y_top = line_bbox[1]
                y_bottom = line_bbox[3]

                # Skip footer/header
                if self._is_footer_or_header(line_text, y_top, y_bottom, page_height):
                    continue

                # Convert PDF coordinates to EMU
                left = int(line_bbox[0] * self.PT_TO_EMU)
                top = int(line_bbox[1] * self.PT_TO_EMU)
                width = int((line_bbox[2] - line_bbox[0]) * self.PT_TO_EMU)
                height = int((line_bbox[3] - line_bbox[1]) * self.PT_TO_EMU)

                # Add padding to cover original text
                pad_h = int(max(font_size, 8) * 0.2 * self.PT_TO_EMU)
                pad_v = int(max(font_size, 8) * 0.15 * self.PT_TO_EMU)
                left = max(0, left - pad_h)
                top = max(0, top - pad_v)
                width = min(width + pad_h * 3, slide_width_emu - left)
                height = min(height + pad_v * 2, slide_height_emu - top)

                # Ensure minimum size
                if width < 50000 or height < 30000:
                    continue

                # Sample background color
                bg_color = self._sample_bg_color(pix, line_bbox, zoom)

                # Get font color
                font_color = self._get_font_color(line)

                # Create text box
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = False
                tf.margin_left = Emu(0)
                tf.margin_right = Emu(0)
                tf.margin_top = Emu(0)
                tf.margin_bottom = Emu(0)

                p = tf.paragraphs[0]
                p.space_before = Pt(0)
                p.space_after = Pt(0)
                run = p.add_run()
                run.text = line_text

                # Font size: use extracted size, default to 10pt for OCR
                effective_size = font_size if font_size > 3 else 10
                run.font.size = Pt(effective_size)
                if is_bold:
                    run.font.bold = True
                if font_color:
                    run.font.color.rgb = font_color

                # Set fill to match surrounding background
                txBox.fill.solid()
                txBox.fill.fore_color.rgb = bg_color

                count += 1

        return count

    def convert_to_pptx(self, input_stream, output_stream, progress_callback=None):
        """Convert PDF to PPTX: each page becomes a slide with image background + text boxes.

        For pages with few extractable text blocks, OCR is used as fallback.
        Returns (output_stream, info_messages).
        """
        try:
            pdf_doc = fitz.open(stream=input_stream.read(), filetype="pdf")
            total_pages = len(pdf_doc)

            first_page = pdf_doc[0]
            pdf_width = first_page.rect.width
            pdf_height = first_page.rect.height

            prs = Presentation()
            slide_width_emu = int(pdf_width * self.PT_TO_EMU)
            slide_height_emu = int(pdf_height * self.PT_TO_EMU)
            prs.slide_width = slide_width_emu
            prs.slide_height = slide_height_emu

            blank_layout = prs.slide_layouts[6]

            info_messages = []
            text_block_count = 0
            ocr_page_count = 0

            zoom = 1.5

            for page_num in range(total_pages):
                page = pdf_doc[page_num]
                page_height = page.rect.height

                slide = prs.slides.add_slide(blank_layout)

                # Render page as image
                mat = fitz.Matrix(zoom, zoom)
                pix = page.get_pixmap(matrix=mat, alpha=False)
                img_bytes = pix.tobytes("jpeg")

                # Add image as background
                img_stream = io.BytesIO(img_bytes)
                bg_pic = slide.shapes.add_picture(
                    img_stream, 0, 0, slide_width_emu, slide_height_emu
                )
                sp = bg_pic._element
                sp.getparent().remove(sp)
                slide.shapes._spTree.insert(2, sp)

                # Try normal text extraction first
                text_dict = page.get_text("dict")
                real_blocks = self._count_real_text_blocks(text_dict, page_height)

                if real_blocks >= self.OCR_THRESHOLD:
                    # Normal extraction - enough text found
                    count = self._add_text_boxes_from_dict(
                        slide, text_dict, page_height, pix, zoom,
                        slide_width_emu, slide_height_emu
                    )
                    text_block_count += count
                elif self._ocr_available:
                    # OCR fallback - page has mostly image-embedded text
                    try:
                        tp = page.get_textpage_ocr(language="eng", dpi=300, full=True)
                        ocr_dict = page.get_text("dict", textpage=tp)
                        count = self._add_text_boxes_from_dict(
                            slide, ocr_dict, page_height, pix, zoom,
                            slide_width_emu, slide_height_emu
                        )
                        text_block_count += count
                        if count > 0:
                            ocr_page_count += 1
                    except Exception as e:
                        # OCR failed for this page, use whatever normal extraction found
                        count = self._add_text_boxes_from_dict(
                            slide, text_dict, page_height, pix, zoom,
                            slide_width_emu, slide_height_emu
                        )
                        text_block_count += count
                else:
                    # No OCR available - use whatever normal extraction found
                    count = self._add_text_boxes_from_dict(
                        slide, text_dict, page_height, pix, zoom,
                        slide_width_emu, slide_height_emu
                    )
                    text_block_count += count

                if progress_callback:
                    progress_callback(min((page_num + 1) / total_pages, 1.0))

            info_messages.append(f"총 {total_pages}페이지, {text_block_count}개 텍스트 블록 변환됨")
            if ocr_page_count > 0:
                info_messages.append(f"OCR 적용: {ocr_page_count}페이지 (이미지에서 텍스트 추출)")
            if not self._ocr_available:
                info_messages.append("⚠️ Tesseract OCR 미설치 - 이미지 내 텍스트는 추출 불가")

            prs.save(output_stream)
            pdf_doc.close()
            return output_stream, info_messages

        except Exception as e:
            print(f"Critical error in convert_to_pptx: {e}")
            traceback.print_exc()
            raise e
